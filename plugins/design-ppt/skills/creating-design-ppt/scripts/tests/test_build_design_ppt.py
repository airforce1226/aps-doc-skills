# -*- coding: utf-8 -*-
import base64
import os
import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
import build_design_ppt as b

SAMPLE = (
    '<!DOCTYPE html><html><head></head><body>\n'
    '<section data-label="표지" data-speaker-notes="표지 노트" '
    'style="width:1920px;height:1080px;background:#0b1b3a;">COVER</section>\n'
    '<section data-label="본문" data-speaker-notes="본문 노트" '
    'style="width:1920px;height:1080px;background:#f5f7fa;">BODY</section>\n'
    '</body></html>'
)

PNG_1x1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+M8AAAMBAQDJ/pLvAAAAAElFTkSuQmCC"
)


def test_split_sections_count_and_fields():
    secs = b.split_sections(SAMPLE)
    assert len(secs) == 2
    assert secs[0]["label"] == "표지"
    assert secs[0]["notes"] == "표지 노트"
    assert "COVER" in secs[0]["html"]
    assert secs[0]["html"].lstrip().startswith("<section")
    assert secs[1]["label"] == "본문"


def test_wrap_section_page_is_full_html():
    page = b.wrap_section_page("<section>X</section>", "body{margin:0}")
    assert page.startswith("<!DOCTYPE html>")
    assert "body{margin:0}" in page
    assert "<section>X</section>" in page
    assert 'charset="utf-8"' in page


def test_find_browser_env_override(tmp_path):
    fake = tmp_path / "chrome.exe"
    fake.write_text("x")
    os.environ["DESIGN_PPT_BROWSER"] = str(fake)
    try:
        assert b.find_browser() == str(fake)
    finally:
        del os.environ["DESIGN_PPT_BROWSER"]


def test_assemble_pptx_fullbleed_notes_author(tmp_path):
    png = tmp_path / "s.png"
    png.write_bytes(PNG_1x1)
    out = tmp_path / "sample v1.0.pptx"
    b.assemble_pptx([{"png": str(png), "notes": "노트내용"}], str(out))

    from pptx import Presentation
    prs = Presentation(str(out))
    assert len(prs.slides) == 1
    assert prs.core_properties.author == "IT전략팀"
    assert prs.slides[0].has_notes_slide
    assert "노트내용" in prs.slides[0].notes_slide.notes_text_frame.text
    # 검토 > 언어 > 맞춤법 검사 안 함 : every notes run carries noProof + lang.
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    tf = prs.slides[0].notes_slide.notes_text_frame
    runs = tf._txBody.findall(".//a:r", ns)
    assert runs, "expected at least one notes run"
    for r in runs:
        rPr = r.find("a:rPr", ns)
        assert rPr is not None and rPr.get("noProof") == "1"
        assert rPr.get("lang") == "ko-KR"
    pic = prs.slides[0].shapes[0]
    assert pic.left == 0 and pic.top == 0
    assert pic.width == prs.slide_width
    assert pic.height == prs.slide_height


def _browser_available():
    try:
        b.find_browser()
        return True
    except SystemExit:
        return False


@pytest.mark.skipif(not _browser_available(), reason="no Chrome/Edge installed")
def test_build_end_to_end(tmp_path):
    deck = tmp_path / "deck.html"
    deck.write_text(SAMPLE, encoding="utf-8")
    out = tmp_path / "sample v1.0.pptx"
    b.build(str(deck), str(out), css="body{margin:0}")

    from pptx import Presentation
    prs = Presentation(str(out))
    assert len(prs.slides) == 2
    assert prs.core_properties.author == "IT전략팀"
    assert prs.slides[0].has_notes_slide
    assert "표지 노트" in prs.slides[0].notes_slide.notes_text_frame.text


def test_assign_page_numbers_skips_cover():
    secs = b.split_sections(SAMPLE)  # 표지 + 본문
    assert b.assign_page_numbers(secs) == [None, "1 / 1"]


def test_assign_page_numbers_honors_off_attr():
    html = SAMPLE.replace('data-label="본문"', 'data-label="본문" data-page-number="off"')
    secs = b.split_sections(html)
    # 표지 제외 + off 제외 -> 번호 대상 0개
    assert b.assign_page_numbers(secs) == [None, None]


def test_assign_page_numbers_sequential_total():
    one = '<section style="width:1920px;height:1080px;background:#f5f7fa;">X</section>'
    secs = b.split_sections("<body>" + (one * 4) + "</body>")
    # 첫 장(표지)만 제외, 나머지 3장은 1~3 / 3
    assert b.assign_page_numbers(secs) == [None, "1 / 3", "2 / 3", "3 / 3"]


def test_inject_page_number_inserts_before_close():
    secs = b.split_sections(SAMPLE)
    out = b.inject_page_number(secs[1]["html"], secs[1]["attrs"], "1 / 1")
    assert "1 / 1" in out
    assert 'data-ppt="text"' in out
    assert out.rstrip().endswith("</section>")
    # 번호 div 는 닫는 태그 '앞'에 들어가야 한다.
    assert out.index("1 / 1") < out.rindex("</section>")


def test_inject_page_number_none_is_noop():
    secs = b.split_sections(SAMPLE)
    assert b.inject_page_number(secs[0]["html"], secs[0]["attrs"], None) == secs[0]["html"]


def test_inject_page_number_dark_vs_light_color():
    secs = b.split_sections(SAMPLE)  # [0]=navy 표지, [1]=paper 본문
    dark = b.inject_page_number(secs[0]["html"], secs[0]["attrs"], "1 / 1")
    light = b.inject_page_number(secs[1]["html"], secs[1]["attrs"], "1 / 1")
    assert "#9fb2d4" in dark      # 밝은 글자색 (네이비 위)
    assert "#9aa6b8" in light     # 슬레이트 글자색 (페이퍼 위)


def test_parse_args_defaults_to_image():
    deck, out, mode = b.parse_args(["deck.html", "제목 v1.0.pptx"])
    assert deck == "deck.html" and out == "제목 v1.0.pptx" and mode == "image"


def test_parse_args_native_mode():
    _, _, mode = b.parse_args(["deck.html", "제목 v1.0.pptx", "--mode", "native"])
    assert mode == "native"
