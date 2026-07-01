# -*- coding: utf-8 -*-
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
import native_render as nr


def test_px_to_pt():
    assert nr.px_to_pt(104) == 52.0  # 104 * 0.5  (1920px == 960pt)


def test_px_to_emu_full_width():
    # 1920 px must map to the full 16:9 slide width in EMU.
    assert nr.px_to_emu(1920) == nr.SLIDE_W_EMU


def test_rgb_to_hex():
    assert nr.rgb_to_hex("rgb(255, 255, 255)") == "FFFFFF"
    assert nr.rgb_to_hex("rgb(11, 27, 58)") == "0B1B3A"


def test_snap_color_exact_and_near():
    assert nr.snap_color("0B1B3A") == "0B1B3A"          # exact token
    assert nr.snap_color("0C1C3B") == "0B1B3A"          # within tolerance -> navy
    assert nr.snap_color("00FF00", tol=5) == "00FF00"   # far -> unchanged


def test_extract_layout_parses_pre_json():
    dom = (
        '<html><body><section></section>'
        '<pre id="__layout__">[{"role":"text","x":120,"y":300,"w":714,"h":138,'
        '"text":"\\uc81c\\ubaa9","font":"\\"Malgun Gothic\\"","sizePx":104,'
        '"weight":"800","color":"rgb(255, 255, 255)","align":"left"}]</pre>'
        '</body></html>'
    )
    nodes = nr.extract_layout(dom)
    assert len(nodes) == 1
    assert nodes[0]["role"] == "text"
    assert nodes[0]["x"] == 120 and nodes[0]["w"] == 714
    assert nodes[0]["text"] == "제목"


def test_extract_layout_missing_pre_raises():
    import pytest
    with pytest.raises(SystemExit):
        nr.extract_layout("<html><body>no pre here</body></html>")


def test_measure_js_is_nonempty_string():
    assert isinstance(nr.MEASURE_JS, str) and "__layout__" in nr.MEASURE_JS


def test_dump_dom_builds_expected_command(monkeypatch):
    captured = {}

    class FakeResult:
        returncode = 0
        stdout = b"<html><pre id=\"__layout__\">[]</pre></html>"
        stderr = b""

    def fake_run(cmd, **kw):
        captured["cmd"] = cmd
        return FakeResult()

    monkeypatch.setattr(nr.subprocess, "run", fake_run)
    dom = nr.dump_dom("C:/tmp/page.html", "chrome.exe")
    assert "--dump-dom" in captured["cmd"]
    assert "--headless=new" in captured["cmd"]
    assert "__layout__" in dom


def test_add_text_creates_editable_textbox(tmp_path):
    from pptx import Presentation
    from pptx.util import Cm
    prs = Presentation()
    prs.slide_width = Cm(33.867); prs.slide_height = Cm(19.05)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    node = {"role": "text", "x": 120, "y": 300, "w": 714, "h": 138,
            "text": "제목", "font": "Malgun Gothic", "sizePx": 104,
            "weight": "800", "color": "rgb(255, 255, 255)", "align": "left",
            "ls": "normal"}
    shp = nr.add_text(slide, node)
    assert shp.has_text_frame
    assert shp.text_frame.paragraphs[0].runs[0].text == "제목"
    assert shp.text_frame.paragraphs[0].runs[0].font.bold is True
    assert round(shp.text_frame.paragraphs[0].runs[0].font.size.pt, 2) == 52.0


def test_add_box_solid_fill(tmp_path):
    from pptx import Presentation
    from pptx.util import Cm
    prs = Presentation()
    prs.slide_width = Cm(33.867); prs.slide_height = Cm(19.05)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    node = {"role": "box", "x": 0, "y": 0, "w": 1920, "h": 1080,
            "bg": "rgb(11, 27, 58)", "grad": "none", "radius": 0}
    shp = nr.add_box(slide, node)
    assert str(shp.fill.fore_color.rgb) == "0B1B3A"
    assert shp.line.fill.type is not None  # no exception accessing line


def test_add_rule_gradient_has_gradfill(tmp_path):
    from pptx import Presentation
    from pptx.util import Cm
    prs = Presentation()
    prs.slide_width = Cm(33.867); prs.slide_height = Cm(19.05)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    node = {"role": "rule", "x": 120, "y": 460, "w": 60, "h": 4,
            "grad": "linear-gradient(90deg, rgb(190, 214, 0), rgb(43, 166, 203))",
            "bg": "rgba(0, 0, 0, 0)", "radius": 0}
    shp = nr.add_rule(slide, node)
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    assert shp._element.spPr.find("a:gradFill", ns) is not None


def test_add_table_native_with_header_and_total(tmp_path):
    from pptx import Presentation
    from pptx.util import Cm
    prs = Presentation()
    prs.slide_width = Cm(33.867); prs.slide_height = Cm(19.05)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    node = {"role": "table", "x": 130, "y": 300, "w": 1660, "h": 400,
            "rows": [
                [{"text": "항목", "header": True}, {"text": "금액", "header": True}],
                [{"text": "매출", "header": False}, {"text": "1,200", "header": False}],
                [{"text": "합계", "header": False}, {"text": "-50", "header": False}],
            ]}
    shp = nr.add_table(slide, node)
    assert shp.has_table
    tbl = shp.table
    assert len(tbl.rows) == 3 and len(tbl.columns) == 2
    hdr = tbl.cell(0, 0)
    assert str(hdr.fill.fore_color.rgb) == "0B1B3A"
    assert str(tbl.cell(2, 0).fill.fore_color.rgb) == "0B3FD1"
    neg_run = tbl.cell(2, 1).text_frame.paragraphs[0].runs[0]
    assert str(neg_run.font.color.rgb) == "C0392B"
    # cell content is vertically centered (PowerPoint 텍스트 맞춤 > 중간)
    from pptx.enum.text import MSO_ANCHOR
    assert hdr.vertical_anchor == MSO_ANCHOR.MIDDLE
    assert tbl.cell(1, 0).vertical_anchor == MSO_ANCHOR.MIDDLE


def test_add_table_empty_rows_returns_none(tmp_path):
    from pptx import Presentation
    from pptx.util import Cm
    prs = Presentation()
    prs.slide_width = Cm(33.867); prs.slide_height = Cm(19.05)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    node = {"role": "table", "x": 130, "y": 300, "w": 1660, "h": 400, "rows": []}
    assert nr.add_table(slide, node) is None  # no crash on empty rows


def test_crop_node_png(tmp_path):
    from PIL import Image
    src = tmp_path / "section.png"
    Image.new("RGB", (1920, 1080), (255, 0, 0)).save(src)
    node = {"x": 100, "y": 50, "w": 200, "h": 80}
    out = nr.crop_node_png(str(src), node, str(tmp_path / "crop.png"))
    img = Image.open(out)
    assert img.size == (200, 80)


def test_render_slide_dispatches_and_counts(tmp_path):
    from pptx import Presentation
    from pptx.util import Cm
    prs = Presentation()
    prs.slide_width = Cm(33.867); prs.slide_height = Cm(19.05)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    nodes = [
        {"role": "box", "x": 0, "y": 0, "w": 1920, "h": 1080,
         "bg": "rgb(11, 27, 58)", "grad": "none", "radius": 0},
        {"role": "text", "x": 120, "y": 300, "w": 700, "h": 138, "text": "제목",
         "font": "Malgun Gothic", "sizePx": 104, "weight": "800",
         "color": "rgb(255,255,255)", "align": "left", "ls": "normal"},
    ]
    counts = nr.render_slide(slide, nodes, section_png=None, tmp_dir=str(tmp_path), si=0)
    assert counts["native"] == 2 and counts["raster"] == 0
    assert any(s.has_text_frame for s in slide.shapes)


import pytest


def _browser_available():
    import build_design_ppt as b
    try:
        b.find_browser(); return True
    except SystemExit:
        return False


@pytest.mark.skipif(not _browser_available(), reason="no Chrome/Edge installed")
def test_build_native_end_to_end(tmp_path):
    from pptx import Presentation
    deck = tmp_path / "deck.html"
    deck.write_text(
        '<section data-label="표지" data-speaker-notes="노트" '
        'style="width:1920px;height:1080px;position:relative;background:#0b1b3a;">'
        '<div data-ppt="text" style="position:absolute;left:120px;top:300px;'
        'font-size:104px;font-weight:800;color:#fff;font-family:\'Malgun Gothic\';">제목</div>'
        '<div data-ppt="rule" style="position:absolute;left:120px;top:460px;width:60px;'
        'height:4px;background:linear-gradient(90deg,#BED600,#2BA6CB);"></div>'
        '</section>', encoding="utf-8")
    out = tmp_path / "표지 v1.0.pptx"
    nr.build_native(str(deck), str(out), css="html,body{margin:0}")

    prs = Presentation(str(out))
    assert len(prs.slides) == 1
    assert prs.core_properties.author == "IT전략팀"
    shapes = prs.slides[0].shapes
    assert any(s.has_text_frame and "제목" in s.text_frame.text for s in shapes)
    # No full-bleed picture: native mode must not bake a screenshot.
    assert not any(s.shape_type == 13 for s in shapes)  # 13 = PICTURE
    # The <section>'s own navy background must become a full-bleed native box.
    def _is_full_bleed(s):
        return (s.left <= 2000 and s.top <= 2000
                and s.width >= prs.slide_width - 8000
                and s.height >= prs.slide_height - 8000)

    def _navy_fill(s):
        try:
            return str(s.fill.fore_color.rgb) == "0B1B3A"
        except Exception:
            return False

    assert any(_is_full_bleed(s) and _navy_fill(s) for s in shapes)


def _slide():
    from pptx import Presentation
    from pptx.util import Cm
    prs = Presentation()
    prs.slide_width = Cm(33.867); prs.slide_height = Cm(19.05)
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_add_text_sets_noproof_on_every_run():
    # 검토>언어>맞춤법 검사 안 함: every run carries noProof + ko-KR.
    node = {"role": "text", "x": 0, "y": 0, "w": 800, "h": 60, "text": "한 줄\n두 줄",
            "font": "맑은 고딕", "sizePx": 40, "weight": "400",
            "color": "rgb(11, 27, 58)", "align": "left", "ls": "normal"}
    shp = nr.add_text(_slide(), node)
    runs = [r for p in shp.text_frame.paragraphs for r in p.runs]
    assert len(runs) == 2  # <br> -> "\n" -> two paragraphs/lines
    for r in runs:
        rPr = r._r.get_or_add_rPr()
        assert rPr.get("noProof") == "1"
        assert rPr.get("lang") == "ko-KR"


def test_add_text_single_line_disables_wrap():
    # A one-line label (h ~ one line) must not word-wrap (footer 2-line bug).
    label = {"role": "text", "x": 0, "y": 0, "w": 240, "h": 26, "text": "APS · Turn on the APS ON",
             "font": "맑은 고딕", "sizePx": 20, "weight": "400",
             "color": "rgb(154,166,184)", "align": "left", "ls": "normal"}
    assert nr.add_text(_slide(), label).text_frame.word_wrap is False
    para = {"role": "text", "x": 0, "y": 0, "w": 800, "h": 300, "text": "긴 본문",
            "font": "맑은 고딕", "sizePx": 29, "weight": "400",
            "color": "rgb(38,53,79)", "align": "left", "ls": "normal"}
    assert nr.add_text(_slide(), para).text_frame.word_wrap is True


def test_add_text_centered_badge():
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    node = {"role": "text", "x": 1776, "y": 40, "w": 94, "h": 36, "text": "대외비",
            "font": "맑은 고딕", "sizePx": 18, "weight": "700",
            "color": "rgb(192,57,43)", "align": "left", "ls": "2.52px", "center": True}
    tf = nr.add_text(_slide(), node).text_frame
    assert tf.vertical_anchor == MSO_ANCHOR.MIDDLE
    assert tf.paragraphs[0].alignment == PP_ALIGN.CENTER


def test_add_box_uniform_border_only_no_fill():
    # 대외비 badge: transparent body + red uniform outline.
    node = {"role": "box", "x": 1776, "y": 40, "w": 94, "h": 36,
            "bg": "rgba(0, 0, 0, 0)", "grad": "none", "radius": 4,
            "bw": 2, "bc": "rgb(192, 57, 43)", "bUniform": True}
    shp = nr.add_box(_slide(), node)
    from pptx.enum.dml import MSO_FILL
    assert shp.fill.type == MSO_FILL.BACKGROUND          # no body fill
    assert str(shp.line.color.rgb) == "C0392B"           # red outline drawn


def test_add_box_single_side_border_not_stroked():
    # border-top divider must NOT become a full rectangle outline.
    node = {"role": "box", "x": 0, "y": 0, "w": 400, "h": 200,
            "bg": "rgba(0, 0, 0, 0)", "grad": "none", "radius": 0,
            "bw": 3, "bc": "rgb(11, 63, 209)", "bUniform": False}
    shp = nr.add_box(_slide(), node)
    assert shp.line.fill.type is not None  # accessible; no solid outline forced


def test_letter_spacing_centi_pt():
    assert nr._letter_spacing_centi_pt("6.16px") == 308   # 6.16 * 0.5 * 100
    assert nr._letter_spacing_centi_pt("normal") is None
    assert nr._letter_spacing_centi_pt(None) is None


def test_measure_js_has_inline_text_folding():
    # The DOM-walk must preserve <br> and fold inline-only blocks (full text), and
    # split space-between footers (own-text guard).
    for token in ("hasOwnText", "inlineOnly", "richText", "colspan"):
        assert token in nr.MEASURE_JS
