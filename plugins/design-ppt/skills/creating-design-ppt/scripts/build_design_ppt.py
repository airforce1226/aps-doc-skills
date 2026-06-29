# -*- coding: utf-8 -*-
"""Build a pixel-faithful .pptx from a Claude-Design-style deck.html.

Each <section> (1920x1080) is rendered to a PNG with headless Chrome/Edge,
then placed full-bleed onto a 16:9 slide. data-speaker-notes -> notes pane.
Author metadata is forced to "IT전략팀".

Usage:
    python build_design_ppt.py deck.html "제목 v1.0.pptx"

Assumptions:
  * <section> elements are siblings (not nested) — true for Claude Design decks.
  * data-label / data-speaker-notes values contain no double-quote character.
"""
import os
import re
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm

AUTHOR = "IT전략팀"
W, H = 1920, 1080
SLIDE_W = Cm(33.867)   # 13.333 in  (16:9)
SLIDE_H = Cm(19.05)    # 7.5 in
ASSETS = Path(__file__).resolve().parent.parent / "assets"

SECTION_RE = re.compile(
    r"<section\b(?P<attrs>[^>]*)>(?P<body>.*?)</section>",
    re.DOTALL | re.IGNORECASE,
)


def _attr(attrs, name):
    m = re.search(re.escape(name) + r'\s*=\s*"([^"]*)"', attrs)
    return m.group(1) if m else ""


def split_sections(html):
    out = []
    for m in SECTION_RE.finditer(html):
        attrs = m.group("attrs")
        out.append({
            "label": _attr(attrs, "data-label"),
            "notes": _attr(attrs, "data-speaker-notes"),
            "attrs": attrs,
            "html": m.group(0),
        })
    return out


# --- Page numbers -----------------------------------------------------------
# The deck is assembled by copying archetypes in arbitrary order, so any page
# number baked into a template would be wrong. Instead we stamp a uniform
# "현재 / 전체" number onto every slide at build time, at a fixed corner position
# so it lands in the same spot on every slide regardless of the archetype.
# 표지(첫 슬라이드)는 관례대로 번호에서 제외하며, 번호 매김은 그 다음 슬라이드에서 1부터
# 시작한다. 특정 슬라이드를 빼려면 <section> 에 data-page-number="off" 를 준다.
PAGE_NUMBER_POS = "right:110px; bottom:50px;"
_OPEN_SECTION_RE = re.compile(r"<section\b[^>]*>", re.IGNORECASE)
_CLOSE_SECTION_RE = re.compile(r"</section\s*>", re.IGNORECASE)


def _section_bg_is_dark(attrs):
    """다크 배경(네이비 표지/간지/마무리)인지 — 페이지 번호 글자색을 밝게 쓰기 위함."""
    style = _attr(attrs, "style")
    m = re.search(r"background(?:-color)?\s*:\s*([^;]+)", style)
    val = m.group(1) if m else ""
    hexm = re.search(r"#([0-9a-fA-F]{6})", val)
    if hexm:
        r, g, b = (int(hexm.group(1)[i:i + 2], 16) for i in (0, 2, 4))
    else:
        rgbm = re.search(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", val)
        if not rgbm:
            return False
        r, g, b = (int(rgbm.group(i)) for i in (1, 2, 3))
    return (0.299 * r + 0.587 * g + 0.114 * b) < 110


def assign_page_numbers(sections):
    """섹션별 페이지 번호 문자열(또는 None) 리스트. 첫 슬라이드(표지)와
    data-page-number="off" 인 슬라이드는 None. 번호는 표시 대상 안에서 1부터."""
    numbered = [
        i != 0 and _attr(sec.get("attrs", ""), "data-page-number").lower() != "off"
        for i, sec in enumerate(sections)
    ]
    total = sum(numbered)
    labels, n = [], 0
    for flag in numbered:
        if flag:
            n += 1
            labels.append("%d / %d" % (n, total))
        else:
            labels.append(None)
    return labels


def inject_page_number(section_html, attrs, label):
    """<section> 닫기 직전에 고정 위치 페이지 번호 <div> 를 자식으로 삽입한다.
    data-ppt="text" 로 네이티브 모드에서도 편집 가능한 텍스트 개체가 되게 한다."""
    if not label:
        return section_html
    color = "#9fb2d4" if _section_bg_is_dark(attrs) else "#9aa6b8"
    badge = (
        '<div data-ppt="text" style="position:absolute; %s z-index:6; '
        "font-family:'Malgun Gothic','맑은 고딕',sans-serif; font-size:20px; "
        'font-weight:600; letter-spacing:.04em; color:%s; text-align:right;">%s</div>'
        % (PAGE_NUMBER_POS, color, label)
    )
    m = None
    for m in _CLOSE_SECTION_RE.finditer(section_html):
        pass  # keep the last </section>
    if not m:
        return section_html
    return section_html[:m.start()] + badge + section_html[m.start():]


def wrap_section_page(section_html, css):
    return (
        '<!DOCTYPE html><html><head><meta charset="utf-8">'
        "<style>html,body{margin:0;padding:0;}" + css + "</style></head>"
        "<body>" + section_html + "</body></html>"
    )


BROWSER_CANDIDATES = [
    r"C:\Program Files\Google\Chrome\Application\chrome.exe",
    r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
    r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
]


def find_browser():
    override = os.environ.get("DESIGN_PPT_BROWSER")
    if override:
        if not os.path.exists(override):
            raise SystemExit("DESIGN_PPT_BROWSER 경로를 찾을 수 없습니다: %s" % override)
        return override
    for p in BROWSER_CANDIDATES:
        if os.path.exists(p):
            return p
    raise SystemExit(
        "Chrome/Edge를 찾을 수 없습니다. DESIGN_PPT_BROWSER 환경변수로 실행파일 경로를 지정하세요."
    )


def capture_slide(page_html_path, png_path, browser):
    url = Path(page_html_path).resolve().as_uri()
    cmd = [
        browser, "--headless=new", "--disable-gpu", "--hide-scrollbars",
        "--force-device-scale-factor=1",
        "--run-all-compositor-stages-before-draw",
        "--virtual-time-budget=3000",
        "--window-size=%d,%d" % (W, H),
        "--screenshot=%s" % png_path,
        url,
    ]
    result = subprocess.run(
        cmd, check=False, timeout=120,
        stdout=subprocess.DEVNULL, stderr=subprocess.PIPE,
    )
    if result.returncode != 0 or not os.path.exists(png_path):
        err = result.stderr.decode(errors="replace").strip() if result.stderr else ""
        raise SystemExit(
            "Chrome 렌더 실패 (exit %d): %s" % (result.returncode, err or png_path)
        )


def _set_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = text
    for p in tf.paragraphs:
        for r in p.runs:
            # _r is python-pptx's documented escape hatch to the lxml element;
            # the public API has no language-tagging, so set the attrs directly.
            rPr = r._r.get_or_add_rPr()
            rPr.set("lang", "ko-KR")
            # 검토 > 언어 > 맞춤법 검사 안 함 : disable spell/grammar check (noProof).
            rPr.set("noProof", "1")


def assemble_pptx(slides, out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # index 6 = "Blank" in python-pptx's default template
    for item in slides:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(item["png"], 0, 0, width=SLIDE_W, height=SLIDE_H)
        if item.get("notes"):
            _set_notes(s, item["notes"])
    prs.core_properties.author = AUTHOR
    prs.core_properties.last_modified_by = AUTHOR
    prs.save(out_path)
    return out_path


def build(deck_path, out_path, css=None):
    html = Path(deck_path).read_text(encoding="utf-8")
    if css is None:
        css_file = ASSETS / "base.css"
        css = css_file.read_text(encoding="utf-8") if css_file.exists() else ""
    sections = split_sections(html)
    if not sections:
        raise SystemExit("deck.html에서 <section>을 찾지 못했습니다.")
    browser = find_browser()
    labels = assign_page_numbers(sections)
    slides = []
    with tempfile.TemporaryDirectory() as tmp:
        for i, sec in enumerate(sections):
            page_path = os.path.join(tmp, "slide_%02d.html" % i)
            sec_html = inject_page_number(sec["html"], sec.get("attrs", ""), labels[i])
            Path(page_path).write_text(wrap_section_page(sec_html, css), encoding="utf-8")
            png_path = os.path.join(tmp, "slide_%02d.png" % i)
            capture_slide(page_path, png_path, browser)
            slides.append({"png": png_path, "notes": sec["notes"]})
        assemble_pptx(slides, out_path)
    print("Rendered %d slide(s)." % len(slides))
    return out_path


def parse_args(argv):
    mode = "image"
    rest = []
    i = 0
    while i < len(argv):
        if argv[i] == "--mode":
            if i + 1 >= len(argv):
                raise SystemExit("--mode 뒤에 image 또는 native 값이 필요합니다.")
            mode = argv[i + 1]
            i += 2
            continue
        rest.append(argv[i])
        i += 1
    if len(rest) < 2:
        raise SystemExit('Usage: python build_design_ppt.py deck.html "제목 v1.0.pptx" [--mode image|native]')
    if mode not in ("image", "native"):
        raise SystemExit("--mode 는 image 또는 native 여야 합니다: %s" % mode)
    return rest[0], rest[1], mode


def main():
    if len(sys.argv) < 3:
        print('Usage: python build_design_ppt.py deck.html "제목 v1.0.pptx" [--mode image|native]',
              file=sys.stderr)
        sys.exit(2)
    deck, out, mode = parse_args(sys.argv[1:])
    if mode == "native":
        import native_render
        native_render.build_native(deck, out)
    else:
        build(deck, out)
    print("Saved:", out)


if __name__ == "__main__":
    main()
