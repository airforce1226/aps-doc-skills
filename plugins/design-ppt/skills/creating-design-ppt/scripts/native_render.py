# -*- coding: utf-8 -*-
"""Native-object renderer for design-ppt (--mode native).

Reads a deck.html section's computed layout (via headless Chrome --dump-dom)
and emits editable python-pptx shapes instead of a baked screenshot.
"""
import json
import os
import re
import subprocess
import tempfile
from pathlib import Path

from PIL import Image

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# Standard 16:9 slide = 13.333" x 7.5". Exact EMU per 작업지시서 §3
# (12192000 x 6858000), so px_to_emu(1) == 6350 exactly (12192000 / 1920).
SLIDE_W_EMU = 12192000
SLIDE_H_EMU = 6858000
PX_TO_PT = 0.54  # design canvas px -> slide pt

# Single source of truth for color snapping — mirrors assets/design-tokens.md.
PALETTE = {
    "navy": "0B1B3A", "navy2": "16263F", "blue": "0B3FD1", "blueSoft": "3F6BD6",
    "blueTint": "7FA3FF", "paper": "F5F7FA", "white": "FFFFFF", "slate": "5B6B85",
    "ink": "26354F", "line": "E1E7F0", "line2": "D8DFE9", "gradLime": "BED600",
    "gradCyan": "2BA6CB", "danger": "C0392B", "success": "1F7A47", "softBlue": "E7EDF8",
    "softBlue2": "EEF2FB", "footer": "9AA6B8",
}


def px_to_pt(px):
    return round(float(px) * PX_TO_PT, 2)


def px_to_emu(px):
    return Emu(int(round(float(px) * SLIDE_W_EMU / 1920)))


def rgb_to_hex(css):
    m = re.search(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", css or "")
    if not m:
        return None
    return "".join("%02X" % int(m.group(i)) for i in (1, 2, 3))


def snap_color(hex6, tol=12):
    """Snap an arbitrary RRGGBB to the nearest palette token within tol; else keep."""
    if not hex6:
        return hex6
    target = hex6.upper()
    tr, tg, tb = (int(target[i:i + 2], 16) for i in (0, 2, 4))
    best, best_d = target, None
    for value in PALETTE.values():
        r, g, bb = (int(value[i:i + 2], 16) for i in (0, 2, 4))
        d = abs(r - tr) + abs(g - tg) + abs(bb - tb)
        if best_d is None or d < best_d:
            best, best_d = value, d
    return best if best_d is not None and best_d <= tol * 3 else target


# Injected before </body>. Walks the DOM, classifies each node, emits layout JSON.
MEASURE_JS = r"""
<script>
(function(){
  function firstFont(f){ return (f||'').split(',')[0].replace(/["']/g,'').trim(); }
  function hasElementChildren(el){
    for (var i=0;i<el.children.length;i++){
      var c=el.children[i], s=getComputedStyle(c);
      if (s.display!=='none' && c.getBoundingClientRect().width>0) return true;
    }
    return false;
  }
  function role(el, cs){
    var hint=el.getAttribute('data-ppt'); if(hint) return hint;
    var tag=el.tagName.toLowerCase();
    if(tag==='table') return 'table';
    if(tag==='img'||tag==='svg'||tag==='canvas') return 'raster';
    var r=el.getBoundingClientRect();
    var grad=(cs.backgroundImage||'').indexOf('gradient')>=0;
    if((r.height<=6 || r.width<=6) && (grad || cs.backgroundColor!=='rgba(0, 0, 0, 0)')) return 'rule';
    var txt=(el.textContent||'').trim();
    if(txt && !hasElementChildren(el)) return 'text';
    if(cs.backgroundColor!=='rgba(0, 0, 0, 0)' || cs.borderTopWidth!=='0px'
       || (cs.borderRadius && cs.borderRadius!=='0px')) return 'box';
    return 'skip';
  }
  var out=[], all=document.querySelectorAll('section *');
  for(var i=0;i<all.length;i++){
    var el=all[i], cs=getComputedStyle(el);
    if(cs.display==='none'||cs.visibility==='hidden') continue;
    var ro=role(el,cs); if(ro==='skip') continue;
    var r=el.getBoundingClientRect();
    if(r.width<=0||r.height<=0) continue;
    var node={role:ro, x:r.left, y:r.top, w:r.width, h:r.height,
      color:cs.color, bg:cs.backgroundColor, grad:cs.backgroundImage,
      align:cs.textAlign, font:firstFont(cs.fontFamily), sizePx:parseFloat(cs.fontSize),
      weight:cs.fontWeight, ls:cs.letterSpacing,
      radius:parseFloat(cs.borderTopLeftRadius)||0};
    if(ro==='text'){ node.text=(el.textContent||'').trim(); }
    if(ro==='table'){
      var rows=[];
      el.querySelectorAll('tr').forEach(function(tr){
        var cells=[];
        tr.querySelectorAll('th,td').forEach(function(td){
          var tcs=getComputedStyle(td);
          cells.push({text:(td.textContent||'').trim(), bg:tcs.backgroundColor,
            color:tcs.color, weight:tcs.fontWeight, align:tcs.textAlign,
            header:td.tagName.toLowerCase()==='th'});
        });
        rows.push(cells);
      });
      node.rows=rows;
    }
    out.push(node);
  }
  var pre=document.createElement('pre'); pre.id='__layout__';
  pre.textContent=JSON.stringify(out); document.body.appendChild(pre);
})();
</script>
"""

_PRE_RE = re.compile(r'<pre id="__layout__">(.*?)</pre>', re.DOTALL)


def extract_layout(dom_text):
    m = _PRE_RE.search(dom_text)
    if not m:
        raise SystemExit("native 모드: __layout__ JSON을 dump-dom에서 찾지 못했습니다.")
    raw = m.group(1)
    # dump-dom HTML-escapes &, <, > inside <pre>; undo before JSON parse.
    raw = raw.replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
    return json.loads(raw)


W, H = 1920, 1080


def dump_dom(page_html_path, browser):
    url = Path(page_html_path).resolve().as_uri()
    cmd = [
        browser, "--headless=new", "--disable-gpu", "--hide-scrollbars",
        "--force-device-scale-factor=1",
        "--run-all-compositor-stages-before-draw",
        "--virtual-time-budget=3000",
        "--window-size=%d,%d" % (W, H),
        "--dump-dom", url,
    ]
    result = subprocess.run(cmd, check=False, timeout=120,
                            stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    dom = result.stdout.decode("utf-8", errors="replace") if result.stdout else ""
    if result.returncode != 0 or "__layout__" not in dom:
        err = result.stderr.decode(errors="replace").strip() if result.stderr else ""
        raise SystemExit("native 모드 dump-dom 실패 (exit %d): %s" % (result.returncode, err))
    return dom


_ALIGN = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT,
          "center": PP_ALIGN.CENTER, "justify": PP_ALIGN.JUSTIFY}


def _is_bold(weight):
    try:
        return int(weight) >= 600
    except (TypeError, ValueError):
        return weight == "bold"


def add_text(slide, node):
    box = slide.shapes.add_textbox(px_to_emu(node["x"]), px_to_emu(node["y"]),
                                   px_to_emu(node["w"]), px_to_emu(node["h"]))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = _ALIGN.get(node.get("align"), PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = node.get("text", "")
    run.font.name = node.get("font") or "맑은 고딕"
    run.font.size = Pt(px_to_pt(node["sizePx"]))
    run.font.bold = _is_bold(node.get("weight"))
    hexc = snap_color(rgb_to_hex(node.get("color")))
    if hexc:
        run.font.color.rgb = RGBColor.from_string(hexc)
    return box


_GRAD_STOP_RE = re.compile(r"rgb\(\s*\d+\s*,\s*\d+\s*,\s*\d+\s*\)")


def _add_rect(slide, node, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, px_to_emu(node["x"]), px_to_emu(node["y"]),
                                 px_to_emu(node["w"]), px_to_emu(node["h"]))
    shp.line.fill.background()  # no border by default
    shp.shadow.inherit = False
    return shp


def add_box(slide, node):
    rounded = float(node.get("radius") or 0) >= 6
    shp = _add_rect(slide, node, rounded=rounded)
    hexc = rgb_to_hex(node.get("bg"))
    if hexc:
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(snap_color(hexc))
    else:
        shp.fill.background()
    return shp


def _set_gradient(shp, stops):
    """Inject a horizontal 2+ stop <a:gradFill> into the shape's spPr via lxml."""
    spPr = shp._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        existing = spPr.find(qn(tag))
        if existing is not None:
            spPr.remove(existing)
    grad = spPr.makeelement(qn("a:gradFill"), {})
    lst = grad.makeelement(qn("a:gsLst"), {})
    n = len(stops)
    for i, hexc in enumerate(stops):
        gs = lst.makeelement(qn("a:gs"), {"pos": str(int(i * 100000 / max(n - 1, 1)))})
        clr = gs.makeelement(qn("a:srgbClr"), {"val": hexc})
        gs.append(clr); lst.append(gs)
    grad.append(lst)
    lin = grad.makeelement(qn("a:lin"), {"ang": "0", "scaled": "1"})
    grad.append(lin)
    ln = spPr.find(qn("a:ln"))
    spPr.insert(list(spPr).index(ln) if ln is not None else len(spPr), grad)


def add_rule(slide, node):
    shp = _add_rect(slide, node, rounded=False)
    stops = [rgb_to_hex(s) for s in _GRAD_STOP_RE.findall(node.get("grad") or "")]
    if len(stops) >= 2:
        _set_gradient(shp, stops)
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(
            snap_color(rgb_to_hex(node.get("bg")) or PALETTE["blue"]))
    return shp


SIZE_BODY_PT = round(14 * PX_TO_PT, 2)  # design body 14px token -> pt
_NUM_RE = re.compile(r"^-?[\d,]+(\.\d+)?$")
_CELL_MARGIN = Emu(45720)  # 0.05 in


def _cell_text(cell, text, *, color, bold, align):
    cell.margin_left = cell.margin_right = _CELL_MARGIN
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "맑은 고딕"
    run.font.size = Pt(SIZE_BODY_PT)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def add_table(slide, node):
    rows = node["rows"]
    nrows, ncols = len(rows), max(len(r) for r in rows)
    gf = slide.shapes.add_table(nrows, ncols, px_to_emu(node["x"]), px_to_emu(node["y"]),
                                px_to_emu(node["w"]), px_to_emu(node["h"]))
    tbl = gf.table
    tbl.first_row = False      # we style cells explicitly; disable built-in banding
    tbl.horz_banding = False
    for ri, row in enumerate(rows):
        is_header = ri == 0 or any(c.get("header") for c in row)
        first_text = (row[0].get("text") if row else "") or ""
        is_total = first_text.strip() in ("합계", "계", "소계", "Total")
        for ci in range(ncols):
            cell = tbl.cell(ri, ci)
            data = row[ci] if ci < len(row) else {"text": ""}
            text = data.get("text", "")
            numeric = bool(_NUM_RE.match(text.strip()))
            negative = text.strip().startswith("-")
            if is_header:
                bg, fg, bold = PALETTE["navy"], PALETTE["white"], True
            elif is_total:
                bg, fg, bold = PALETTE["blue"], PALETTE["white"], True
            elif ri % 2 == 0:
                bg, fg, bold = PALETTE["softBlue2"], PALETTE["ink"], False
            else:
                bg, fg, bold = PALETTE["white"], PALETTE["ink"], False
            if negative and not is_header:
                fg = PALETTE["danger"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(bg)
            align = PP_ALIGN.RIGHT if (numeric and not is_header) else PP_ALIGN.LEFT
            _cell_text(cell, text, color=fg, bold=bold, align=align)
    return gf


def crop_node_png(section_png, node, out_path):
    img = Image.open(section_png)
    x, y = int(round(node["x"])), int(round(node["y"]))
    w, h = int(round(node["w"])), int(round(node["h"]))
    img.crop((x, y, x + w, y + h)).save(out_path)
    return out_path


def add_raster(slide, node, section_png, tmp_dir, idx):
    out = "%s/raster_%d.png" % (tmp_dir, idx)
    crop_node_png(section_png, node, out)
    return slide.shapes.add_picture(out, px_to_emu(node["x"]), px_to_emu(node["y"]),
                                    px_to_emu(node["w"]), px_to_emu(node["h"]))


def render_slide(slide, nodes, section_png, tmp_dir, si):
    counts = {"native": 0, "raster": 0}
    # Boxes/rules first (background layer), then raster, then tables/text on top.
    order = {"box": 0, "rule": 1, "raster": 2, "table": 3, "text": 4}
    for node in sorted(nodes, key=lambda n: order.get(n["role"], 5)):
        role = node["role"]
        if role == "box":
            add_box(slide, node); counts["native"] += 1
        elif role == "rule":
            add_rule(slide, node); counts["native"] += 1
        elif role == "text":
            add_text(slide, node); counts["native"] += 1
        elif role == "table":
            add_table(slide, node); counts["native"] += 1
        elif role == "raster":
            if section_png:
                add_raster(slide, node, section_png, tmp_dir, len(slide.shapes))
                counts["raster"] += 1
    return counts


def build_native(deck_path, out_path, css=None, browser=None):
    import build_design_ppt as b
    from pptx import Presentation
    html = Path(deck_path).read_text(encoding="utf-8")
    if css is None:
        css_file = b.ASSETS / "base.css"
        css = css_file.read_text(encoding="utf-8") if css_file.exists() else ""
    sections = b.split_sections(html)
    if not sections:
        raise SystemExit("deck.html에서 <section>을 찾지 못했습니다.")
    browser = browser or b.find_browser()
    prs = Presentation()
    prs.slide_width = b.SLIDE_W
    prs.slide_height = b.SLIDE_H
    blank = prs.slide_layouts[6]
    report = []
    with tempfile.TemporaryDirectory() as tmp:
        for si, sec in enumerate(sections):
            page = os.path.join(tmp, "slide_%02d.html" % si)
            wrapped = b.wrap_section_page(sec["html"], css).replace(
                "</body>", MEASURE_JS + "</body>")
            Path(page).write_text(wrapped, encoding="utf-8")
            nodes = extract_layout(dump_dom(page, browser))
            section_png = None
            if any(n["role"] == "raster" for n in nodes):
                section_png = os.path.join(tmp, "slide_%02d.png" % si)
                b.capture_slide(page, section_png, browser)
            slide = prs.slides.add_slide(blank)
            counts = render_slide(slide, nodes, section_png, tmp, si)
            if sec.get("notes"):
                b._set_notes(slide, sec["notes"])
            report.append((si, sec.get("label", ""), counts))
        prs.core_properties.author = b.AUTHOR
        prs.core_properties.last_modified_by = b.AUTHOR
        prs.save(out_path)
    print_report(report)
    return out_path


def print_report(report):
    tot_n = sum(c["native"] for _, _, c in report)
    tot_r = sum(c["raster"] for _, _, c in report)
    for si, label, c in report:
        print("Slide %02d (%s): native=%d raster=%d"
              % (si + 1, label or "-", c["native"], c["raster"]))
    print("Total: %d slides, %d native objects, %d raster fallbacks."
          % (len(report), tot_n, tot_r))
