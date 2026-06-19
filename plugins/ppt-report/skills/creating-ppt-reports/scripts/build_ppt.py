# -*- coding: utf-8 -*-
"""Build a styled Korean .pptx report from a JSON spec.

Usage:
    python build_ppt.py spec.json "output v1.0.pptx"

This is the PPT half of the docs->ppt workflow. Author it AFTER the source
.docx exists; derive this spec from that document's sections.

Team rules baked in (see the creating-ppt-reports skill):
  * Author metadata is forced to "IT전략팀" (never a person's name).
  * Every run gets lang="ko-KR" and noProof="1" so PowerPoint shows no
    spelling/grammar red squiggles. The in-app final step
    (검토 > 언어 > 교정 언어 설정) still applies as a manual confirm.
  * Output filename: spaces (no underscores) + version suffix "vN.0".
    The caller is responsible for passing such a path.

Spec format (UTF-8 JSON):
{
  "title": "보고서 제목",
  "subtitle": "부제 (선택)",
  "meta": "작성일 2026-06-19 | 작성 IT전략팀  (선택)",
  "accent": "1F6F43",            # 강조색 hex (선택, 기본 녹색)
  "font": "맑은 고딕",            # 글꼴 (선택)
  "slides": [
    {
      "heading": "1. 사업 개요",
      "blocks": [
        {"bullet": "내용", "level": 0, "label": "사업명"},   # label/level 선택
        {"para": "보조 설명", "grey": true},                  # grey 선택
        {"table": {
            "headers": ["열1", "열2"],
            "rows": [["a", "b"], ["c", "TBD"]]
        }}
      ]
    }
  ],
  "footnote": "※ ... (선택, 마지막 슬라이드 하단)"
}

값이 "TBD" 인 셀/문구는 회색 이탤릭 "(미정 — 추후 확정)" 으로 렌더링된다.
"""
import json
import sys

from pptx import Presentation
from pptx.util import Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

GREY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TBD = "TBD"
AUTHOR = "IT전략팀"

# 16:9 slide is 13.333 x 7.5 inches
SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)
MARGIN = Cm(1.6)


def _no_proof(run, font_name):
    """Set Korean language + no-proof on a run so red squiggles never show."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "ko-KR")
    rPr.set("noProof", "1")
    # east-asian font binding
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        latin = rPr.makeelement(qn("a:latin"), {})
        rPr.append(latin)
    latin.set("typeface", font_name)
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font_name)


def build(spec, out_path):
    accent = RGBColor.from_string(spec.get("accent", "1F6F43"))
    font_name = spec.get("font", "맑은 고딕")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    def add_run(p, text, *, bold=False, size=14, color=None, italic=False):
        r = p.add_run()
        r.text = text
        r.font.bold = bold
        r.font.italic = italic
        r.font.size = Pt(size)
        r.font.name = font_name
        if color is not None:
            r.font.color.rgb = color
        _no_proof(r, font_name)
        return r

    def textbox(slide, left, top, width, height):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        return tb, tf

    # ---------- title slide ----------
    s = prs.slides.add_slide(blank)
    tb, tf = textbox(s, MARGIN, Cm(6.5), SLIDE_W - 2 * MARGIN, Cm(6))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, spec.get("title", ""), bold=True, size=36, color=accent)
    if spec.get("subtitle"):
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(14)
        add_run(p, spec["subtitle"], size=18, color=GREY)
    if spec.get("meta"):
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(24)
        add_run(p, spec["meta"], size=13, color=GREY)

    # ---------- content slides ----------
    slides = spec.get("slides", [])
    for idx, sl in enumerate(slides):
        s = prs.slides.add_slide(blank)
        # heading
        tb, tf = textbox(s, MARGIN, Cm(0.9), SLIDE_W - 2 * MARGIN, Cm(2))
        p = tf.paragraphs[0]
        add_run(p, sl.get("heading", ""), bold=True, size=24, color=accent)
        # accent underline bar
        bar = s.shapes.add_shape(
            1, MARGIN, Cm(2.7), SLIDE_W - 2 * MARGIN, Pt(2.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        bar.shadow.inherit = False

        # body
        body_top = Cm(3.3)
        tb, tf = textbox(
            s, MARGIN, body_top, SLIDE_W - 2 * MARGIN, SLIDE_H - body_top - Cm(1.2)
        )
        first = True
        last_table = None
        for block in sl.get("blocks", []):
            if "bullet" in block:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.level = block.get("level", 0)
                p.space_after = Pt(6)
                add_run(p, "• ", bold=True, size=16, color=accent)
                if block.get("label"):
                    add_run(p, block["label"] + " : ", bold=True, size=16)
                add_run(p, str(block["bullet"]), size=16)
            elif "para" in block:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.space_after = Pt(6)
                grey = block.get("grey", False)
                add_run(p, str(block["para"]), size=14,
                        color=GREY if grey else None, italic=grey)
            elif "table" in block:
                last_table = block["table"]

        # a table block, if any, is drawn as a real table below the text
        if last_table is not None:
            _add_table(s, last_table, accent, font_name, add_run)

    # ---------- footnote on last slide ----------
    if spec.get("footnote") and prs.slides:
        s = prs.slides[-1]
        tb, tf = textbox(s, MARGIN, SLIDE_H - Cm(1.4),
                         SLIDE_W - 2 * MARGIN, Cm(1.0))
        p = tf.paragraphs[0]
        add_run(p, spec["footnote"], size=10, color=GREY, italic=True)

    prs.core_properties.author = AUTHOR
    prs.core_properties.last_modified_by = AUTHOR
    prs.save(out_path)
    return out_path


def _add_table(slide, spec_t, accent, font_name, add_run):
    headers = spec_t["headers"]
    rows = spec_t.get("rows", [])
    ncols = len(headers)
    nrows = len(rows) + 1
    left, top = MARGIN, Cm(9.5)
    width = SLIDE_W - 2 * MARGIN
    height = Cm(0.9) * nrows
    gtable = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
    for c, htext in enumerate(headers):
        cell = gtable.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        add_run(p, str(htext), bold=True, size=12, color=WHITE)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gtable.cell(r, c)
            text = "" if val is None else str(val)
            is_tbd = text.strip() == TBD
            p = cell.text_frame.paragraphs[0]
            add_run(p, "(미정 — 추후 확정)" if is_tbd else text, size=11,
                    color=GREY if is_tbd else None, italic=is_tbd)


def main():
    if len(sys.argv) < 3:
        print('Usage: python build_ppt.py spec.json "output v1.0.pptx"',
              file=sys.stderr)
        sys.exit(2)
    with open(sys.argv[1], "r", encoding="utf-8") as f:
        spec = json.load(f)
    out = build(spec, sys.argv[2])
    print("Saved:", out)


if __name__ == "__main__":
    main()
